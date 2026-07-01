"""
Purpose: One copyable tool that reads any file — text, images, PDF, LaTeX, PPTX, Word, audio, video — and returns the raw extracted content for the agent.
LLM-Note:
  Dependencies: imports from [base64, mimetypes, pathlib, shutil, subprocess, tempfile, charset_normalizer, docx, pypdf, pptx, connectonion.transcribe] | requires ffmpeg on PATH for video | imported by [copied into user projects via `co copy read_file`; not re-exported from useful_tools/__init__]
  Data flow: read_file(path) → dispatch by extension → text/PDF/PPTX/DOCX return extracted text | images return a data:image/...;base64 URL (picked up by the image_result_formatter plugin so the vision model sees it) | audio transcribes via transcribe() | video extracts audio with ffmpeg then transcribes
  State/Effects: reads the file from disk | audio/video make a transcription API call (transcribe) | video shells out to ffmpeg into a temp dir | no caching
  Integration: exposes read_file(path) as an agent tool | supersedes the built-in text read_file — pass this one, not both (same tool name) | image results rely on the image_result_formatter plugin being enabled
  Errors: returns "Error: ..." strings for the tool's own boundary checks (missing file, unsupported binary, legacy .ppt/.doc, ffmpeg missing) | lets transcribe()/pypdf/pptx/docx exceptions bubble to the tool executor

Read any file and return its contents for the agent to use.

Dispatches by extension:
  - Text/code/markup (.txt .md .csv .json .tex .py …) → raw text
  - Images (.png .jpg .jpeg .gif .webp)              → data URL for the vision model
  - PDF (.pdf)                                       → extracted text, per page
  - PowerPoint (.pptx)                               → extracted slide text
  - Word (.docx)                                     → extracted document text
  - Audio (.mp3 .wav .m4a …)                         → transcript
  - Video (.mp4 .mov .mkv …)                         → audio transcript (via ffmpeg)

This supersedes the built-in `read_file` (text + line numbers). When you copy it
in, pass this one to your agent — not both (they share the tool name `read_file`).

Usage:
    from connectonion import Agent
    from tools.read_file import read_file

    agent = Agent("assistant", tools=[read_file])

Video needs `ffmpeg` on PATH; the tool says so if it's missing.
"""

import base64
import mimetypes
import shutil
import subprocess
import tempfile
from pathlib import Path

from charset_normalizer import from_bytes
from connectonion import transcribe
from docx import Document
from pypdf import PdfReader
from pptx import Presentation

# Extensions handled by a dedicated path; everything else is read as text.
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
AUDIO_EXTS = {".mp3", ".wav", ".aiff", ".aac", ".ogg", ".flac", ".m4a", ".webm"}
VIDEO_EXTS = {".mp4", ".mov", ".avi", ".mkv", ".flv", ".wmv", ".m4v"}


def read_file(path: str) -> str:
    """Read any file and hand you its contents. Dispatches by extension:
    - text / code / .md / .csv / .json / .tex  →  the file's text, as-is
    - images (.png .jpg .jpeg .gif .webp)  →  the image enters your vision; after reading you can SEE and describe it
    - PDF (.pdf)  →  the document's text, page by page
    - PowerPoint (.pptx)  →  the text of each slide
    - Word (.docx)  →  the document's text
    - audio (.mp3 .wav …) / video (.mp4 …)  →  a transcript of the speech
    Call this whenever you are given a file path — including images — before ever saying you cannot read or see it.

    Args:
        path: Path to the file to read.

    Returns:
        Extracted text for most types. Images return a `data:image/...;base64,...`
        data URL that the image_result_formatter plugin turns into an image the
        model can see. On a boundary problem (missing file, unsupported binary,
        legacy .ppt/.doc, missing ffmpeg) returns an "Error: ..." message you can act on.
    """
    file_path = Path(path)
    if not file_path.exists():
        return f"Error: File '{path}' does not exist"
    if not file_path.is_file():
        return f"Error: '{path}' is not a file"

    ext = file_path.suffix.lower()

    if ext in IMAGE_EXTS:
        return _read_image(file_path)
    if ext == ".pdf":
        return _read_pdf(file_path)
    if ext in (".pptx", ".ppt"):
        return _read_pptx(file_path, ext)
    if ext in (".docx", ".doc"):
        return _read_docx(file_path, ext)
    if ext in AUDIO_EXTS:
        return transcribe(str(file_path))
    if ext in VIDEO_EXTS:
        return _read_video(file_path)
    return _read_text(file_path, ext)


def _read_text(file_path: Path, ext: str) -> str:
    """Read a text/code/markup file. Rejects binary files."""
    data = file_path.read_bytes()
    # NUL bytes almost never appear in text; their presence means binary.
    if b"\x00" in data[:8192]:
        return (
            f"Error: '{file_path}' looks like a binary file with an unsupported "
            f"extension ('{ext or 'none'}'), not text."
        )
    # Detect the encoding, limited to UTF-8 and Chinese GB18030 so short snippets
    # aren't misdetected as some unrelated charset.
    match = from_bytes(data, cp_isolation=["utf_8", "gb18030"]).best()
    return str(match) if match is not None else data.decode("utf-8", errors="replace")


def _read_image(file_path: Path) -> str:
    """Return an image as a base64 data URL for the vision model."""
    mime = mimetypes.guess_type(str(file_path))[0] or "image/png"
    b64 = base64.b64encode(file_path.read_bytes()).decode("utf-8")
    return f"data:{mime};base64,{b64}"


def _read_pdf(file_path: Path) -> str:
    """Extract text from a PDF, one section per page."""
    texts = [(page.extract_text() or "").strip() for page in PdfReader(str(file_path)).pages]
    if not any(texts):
        return (
            f"Error: no extractable text in '{file_path}'. "
            "It may be a scanned PDF (images only)."
        )
    return "\n\n".join(f"--- Page {i} ---\n{t}" for i, t in enumerate(texts, 1))


def _read_pptx(file_path: Path, ext: str) -> str:
    """Extract slide text from a .pptx presentation."""
    if ext == ".ppt":
        return (
            "Error: legacy .ppt is not supported. Convert it to .pptx "
            "(e.g. with LibreOffice or PowerPoint) and try again."
        )
    slides = []
    for i, slide in enumerate(Presentation(str(file_path)).slides, 1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    return "\n\n".join(slides).strip()


def _read_docx(file_path: Path, ext: str) -> str:
    """Extract text from a .docx document, paragraph by paragraph."""
    if ext == ".doc":
        return (
            "Error: legacy .doc is not supported. Convert it to .docx "
            "(e.g. with LibreOffice or Word) and try again."
        )
    paras = [p.text for p in Document(str(file_path)).paragraphs if p.text.strip()]
    return "\n".join(paras)


def _read_video(file_path: Path) -> str:
    """Extract the audio track with ffmpeg, then transcribe it."""
    if shutil.which("ffmpeg") is None:
        return (
            "Error: reading video requires ffmpeg on PATH. "
            "Install it (e.g. 'brew install ffmpeg' or 'apt install ffmpeg') and try again."
        )
    with tempfile.TemporaryDirectory() as tmp:
        audio_path = Path(tmp) / "audio.mp3"
        proc = subprocess.run(
            ["ffmpeg", "-y", "-i", str(file_path),
             "-vn", "-ac", "1", "-b:a", "64k", str(audio_path)],
            capture_output=True, text=True,
        )
        if proc.returncode != 0 or not audio_path.exists():
            return f"Error: failed to extract audio from '{file_path}':\n{proc.stderr[-500:]}"
        return transcribe(str(audio_path))
