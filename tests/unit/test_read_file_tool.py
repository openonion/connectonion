"""Unit tests for connectonion/useful_tools/read_file.py (the `co copy read_file` tool).

Covers the local dispatch paths and boundary errors. pypdf and python-pptx are
real dependencies, so PDF/PPTX extraction is tested against real files. Audio/video
transcription is exercised only through its dispatch + boundary branches (no
network/API): audio dispatch calls transcribe(); video without ffmpeg and ffmpeg
failure return actionable errors.
"""

import base64
import shutil
import pytest
from pathlib import Path
from unittest.mock import Mock
from docx import Document
from pptx import Presentation

import connectonion.useful_tools.read_file as read_file_module
from connectonion.useful_tools.read_file import read_file
from connectonion.useful_plugins.image_result_formatter import _format_image_result


# A minimal PDF whose single page contains the text "Hello PDF World".
_MINIMAL_PDF = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 200 200]/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj
4 0 obj<</Length 44>>stream
BT /F1 18 Tf 20 100 Td (Hello PDF World) Tj ET
endstream endobj
5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj
xref
0 6
0000000000 65535 f
trailer<</Root 1 0 R/Size 6>>
startxref
0
%%EOF"""


class TestBoundaries:
    def test_missing_file(self, tmp_path):
        result = read_file(str(tmp_path / "nope.txt"))
        assert result == f"Error: File '{tmp_path / 'nope.txt'}' does not exist"

    def test_directory_is_not_a_file(self, tmp_path):
        result = read_file(str(tmp_path))
        assert "is not a file" in result


class TestText:
    def test_plain_text_returned_as_is(self, tmp_path):
        f = tmp_path / "note.txt"
        f.write_text("hello\nworld\n", encoding="utf-8")
        # Raw content, no line numbers and no truncation (unlike the built-in read_file).
        assert read_file(str(f)) == "hello\nworld\n"

    def test_latex_source_is_read_as_text(self, tmp_path):
        f = tmp_path / "paper.tex"
        f.write_text(r"\documentclass{article}\begin{document}Hi\end{document}", encoding="utf-8")
        assert r"\documentclass{article}" in read_file(str(f))

    def test_json_and_code_read_as_text(self, tmp_path):
        f = tmp_path / "data.json"
        f.write_text('{"a": 1}', encoding="utf-8")
        assert read_file(str(f)) == '{"a": 1}'

    def test_non_utf8_chinese_is_decoded(self, tmp_path):
        # GBK-encoded Chinese (common on Windows) must not come back as mojibake.
        f = tmp_path / "gbk.txt"
        f.write_bytes("你好，世界。这是一段中文。".encode("gbk"))
        assert read_file(str(f)) == "你好，世界。这是一段中文。"

    def test_large_text_returned_whole(self, tmp_path):
        # Contract (#169): return raw content, no pre-formatting / truncation.
        f = tmp_path / "big.txt"
        body = "a" * 250_000
        f.write_text(body, encoding="utf-8")
        assert read_file(str(f)) == body

    def test_unknown_binary_extension_rejected(self, tmp_path):
        f = tmp_path / "blob.dat"
        f.write_bytes(b"MZ\x00\x00\x90\x00binary\x00stuff")
        result = read_file(str(f))
        assert result.startswith("Error:")
        assert "binary" in result


class TestImage:
    def test_png_returns_data_url(self, tmp_path):
        f = tmp_path / "pic.png"
        payload = b"\x89PNG\r\n\x1a\nfake-image-bytes"
        f.write_bytes(payload)
        result = read_file(str(f))
        assert result.startswith("data:image/png;base64,")
        # Round-trips to the original bytes so the vision model sees the real image.
        b64 = result.split(",", 1)[1]
        assert base64.b64decode(b64) == payload

    def test_jpg_mime(self, tmp_path):
        f = tmp_path / "pic.jpg"
        f.write_bytes(b"\xff\xd8\xff\xe0jpeg")
        assert read_file(str(f)).startswith("data:image/jpeg;base64,")


class TestImageReachesVisionModel:
    """read_file's image output must actually arrive at the model as an image.

    The tool returns a data URL (data only); the image_result_formatter plugin
    turns it into an image_url message the vision model can see — same path as
    browser screenshots. This locks that end-to-end contract.
    """

    def test_image_output_becomes_image_url_message(self, tmp_path):
        f = tmp_path / "shot.png"
        f.write_bytes(b"\x89PNG\r\n\x1a\nfake-image-bytes")
        data_url = read_file(str(f))
        assert data_url.startswith("data:image/png;base64,")

        # Simulate the tool result sitting in a session, then run the plugin.
        agent = Mock()
        agent.io = None
        agent.current_session = {
            "trace": [{
                "type": "tool_result", "status": "success",
                "name": "read_file", "tool_id": "call_1", "result": data_url,
            }],
            "messages": [{
                "role": "tool", "tool_call_id": "call_1", "content": data_url,
            }],
        }

        _format_image_result(agent)

        messages = agent.current_session["messages"]
        assert len(messages) == 2
        image_msg = messages[1]
        assert image_msg["role"] == "user"
        image_part = next(p for p in image_msg["content"] if p["type"] == "image_url")
        # The whole image reaches the model — not a truncated/altered blob.
        assert image_part["image_url"]["url"] == data_url


class TestPdfPptx:
    def test_pdf_extracts_text(self, tmp_path):
        f = tmp_path / "doc.pdf"
        f.write_bytes(_MINIMAL_PDF)
        result = read_file(str(f))
        assert "Hello PDF World" in result
        assert "--- Page 1 ---" in result

    def test_pptx_extracts_slide_text(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Quarterly Results"
        f = tmp_path / "deck.pptx"
        prs.save(str(f))
        result = read_file(str(f))
        assert "--- Slide 1 ---" in result
        assert "Quarterly Results" in result

    def test_legacy_ppt_rejected(self, tmp_path):
        f = tmp_path / "old.ppt"
        f.write_bytes(b"\xd0\xcf\x11\xe0 legacy ppt")
        result = read_file(str(f))
        assert "legacy .ppt is not supported" in result

    def test_docx_extracts_text(self, tmp_path):
        doc = Document()
        doc.add_paragraph("Contract Agreement")
        doc.add_paragraph("Party A shall deliver the goods.")
        f = tmp_path / "contract.docx"
        doc.save(str(f))
        result = read_file(str(f))
        assert "Contract Agreement" in result
        assert "Party A shall deliver the goods." in result

    def test_legacy_doc_rejected(self, tmp_path):
        f = tmp_path / "old.doc"
        f.write_bytes(b"\xd0\xcf\x11\xe0 legacy doc")
        result = read_file(str(f))
        assert "legacy .doc is not supported" in result


class TestAudioVideoDispatch:
    def test_audio_dispatches_to_transcribe(self, tmp_path, monkeypatch):
        calls = {}

        def fake_transcribe(path):
            calls["path"] = path
            return "TRANSCRIPT"

        monkeypatch.setattr(read_file_module, "transcribe", fake_transcribe)
        f = tmp_path / "clip.mp3"
        f.write_bytes(b"ID3 fake audio")
        assert read_file(str(f)) == "TRANSCRIPT"
        assert calls["path"] == str(f)

    def test_video_without_ffmpeg(self, tmp_path, monkeypatch):
        monkeypatch.setattr(read_file_module.shutil, "which", lambda name: None)
        f = tmp_path / "clip.mp4"
        f.write_bytes(b"\x00\x00\x00\x18ftyp fake video")
        result = read_file(str(f))
        assert result.startswith("Error:")
        assert "ffmpeg" in result

    @pytest.mark.skipif(shutil.which("ffmpeg") is None, reason="ffmpeg not installed")
    def test_video_ffmpeg_failure_reports_error(self, tmp_path):
        # Real ffmpeg on non-video bytes exits non-zero -> actionable error, no transcription.
        f = tmp_path / "clip.mp4"
        f.write_bytes(b"not actually a video")
        result = read_file(str(f))
        assert result.startswith("Error: failed to extract audio")
